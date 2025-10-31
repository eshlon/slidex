from fastapi.middleware.cors import CORSMiddleware
import os
import io
import json
import logging
import subprocess
import tempfile
import platform
from dotenv import load_dotenv
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from openai import OpenAI
import requests

from fastapi import FastAPI, HTTPException
from pydantic import BaseModel
from starlette.responses import StreamingResponse

# Load environment variables from .env file
load_dotenv()

# Configure logging
logging.basicConfig(
    format="%(asctime)s - %(levelname)s - %(message)s", level=logging.INFO
)
logger = logging.getLogger(__name__)

# Initialize OpenAI client
client = OpenAI(api_key=os.environ.get("OPENAI_API_KEY", ""))

# Get Pexels API key and template paths
PEXELS_API_KEY = os.environ.get("PEXELS_API_KEY", "")
DEFAULT_TEMPLATE_PATH = os.environ.get("TEMPLATE_PATH", "pitch.pptx")

# Template mapping
TEMPLATE_PATHS = {
    "aura":                         "public/templates/aura.pptx",
    "bevel-design":                 "public/templates/bevel-design.pptx",
    "big-bold":                     "public/templates/big-bold.pptx",
    "blue-spheres":                 "public/templates/blue-spheres.pptx",
    "bold-underlines":              "public/templates/bold-underlines.pptx",
    "color-block":                  "public/templates/color-block.pptx",
    "corporate-blue":               "public/templates/corporate-blue.pptx",
    "fission":                      "public/templates/fission.pptx",
    "futuristic-pitch":             "public/templates/futuristic-pitch.pptx",
    "gradient-aura":                "public/templates/gradient-aura.pptx",
    "luminous-design":              "public/templates/luminous-design.pptx",
    "modern-minimal":               "public/templates/modern-minimal.pptx",
    "modern-sales-strategy":        "public/templates/modern-sales-strategy.pptx",
    "oasis-design":                 "public/templates/oasis-design.pptx",
    "pantone-2022":                 "public/templates/pantone-2022.pptx",
    "simple-budget-planning":       "public/templates/simple-budget-planning.pptx",
    "sleek-corporate-finance":      "public/templates/sleek-corporate-finance.pptx",
    "vibrant-creative":             "public/templates/vibrant-creative.pptx",
}

def get_template_structure(template_path):
    """Extracts a brief structure of layouts and placeholders from a PowerPoint template, including image placeholders."""
    prs = Presentation(template_path)
    structure = {"masters": []}

    placeholder_map = {
        PP_PLACEHOLDER.TITLE: "Title",
        PP_PLACEHOLDER.BODY: "Body",
        PP_PLACEHOLDER.CENTER_TITLE: "Center Title",
        PP_PLACEHOLDER.SUBTITLE: "Subtitle",
        PP_PLACEHOLDER.DATE: "Date",
        PP_PLACEHOLDER.FOOTER: "Footer",
        PP_PLACEHOLDER.HEADER: "Header",
        PP_PLACEHOLDER.SLIDE_NUMBER: "Slide Number",
        PP_PLACEHOLDER.PICTURE: "Picture",
    }

    for master_idx, master in enumerate(prs.slide_masters):
        master_data = {"master_idx": master_idx, "layouts": []}
        for layout_idx, layout in enumerate(master.slide_layouts):
            layout_data = {"layout_idx": layout_idx, "name": layout.name, "placeholders": []}
            for shape in layout.shapes:
                if shape.is_placeholder:
                    placeholder_type = shape.placeholder_format.type
                    ph_type_name = placeholder_map.get(placeholder_type, "Unknown")
                    if ph_type_name in ["Title", "Body", "Center Title", "Subtitle", "Date", "Footer", "Header", "Slide Number", "Picture"]:
                        placeholder_data = {
                            "idx": shape.placeholder_format.idx,
                            "id": shape.placeholder_format.idx,
                            "name": shape.name,
                            "type": ph_type_name
                        }
                        layout_data["placeholders"].append(placeholder_data)
            master_data["layouts"].append(layout_data)
        structure["masters"].append(master_data)
    return structure

def get_pexels_image_url(query, api_key):
    """Fetches an image URL from Pexels based on the search query."""
    if not api_key:
        return None
    url = "https://api.pexels.com/v1/search"
    headers = {"Authorization": api_key}
    params = {"query": query, "per_page": 1}
    try:
        response = requests.get(url, headers=headers, params=params, timeout=10)
        if response.status_code == 200:
            data = response.json()
            if data["photos"]:
                return data["photos"][0]["src"]["large"]
    except Exception as e:
        logger.error(f"Error fetching from Pexels: {e}")
    return None

def generate_content_outline(user_content):
    """Generates a presentation content outline using OpenAI GPT based on user content."""
    sys_prompt = """
    You are Slidex, an expert AI assistant that generate high-quality content based on the provided content by user or user's query.
    your task is to generate, curate, edit, and finalize textual material for a excellent presentation.
    If the user provided you with prepaid content, try to edit or finalize the content as user asker.
    If user asked about some topic generate content.
    repsond only with final content, as a presentation outline.
    """
    
    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": sys_prompt},
                {"role": "user", "content": user_content}
            ],
            timeout=60
        )
        return response.choices[0].message.content
    except Exception as e:
        logger.error(f"Error generating content outline: {str(e)}")
        return None

def generate_structured_outline(prompt: str, slide_count: int, language: str):
    """Generates a structured presentation outline in JSON format using an AI model."""
    
    system_prompt = f"""
You are Slidex, an expert AI assistant that creates high-quality presentation outlines.

Your task is to generate a structured presentation outline based on the user's topic, desired number of slides, and language preference. The goal is to produce rich, detailed content for each slide.

Requirements:
- Create exactly {slide_count} slides.
- Write in {language}.
- Each slide must have a "title" (string) and "content" (an array of 4-6 detailed strings).
- Each item in the "content" array should be a full, descriptive sentence or a comprehensive bullet point, providing substantial information. Avoid one or two-word points.
- The first slide should be an introduction or title slide.
- The last slide should be a conclusion or summary slide.
- The middle slides should logically flow from one to the next, building a comprehensive narrative.

You MUST respond with ONLY a valid JSON object containing a single key "slides" which is an array of objects, without any surrounding text, comments, or markdown. The JSON should strictly follow this format:
{{
  "slides": [
    {{
      "title": "Slide Title",
      "content": ["This is the first detailed bullet point.", "This is the second detailed bullet point, explaining more about the topic."]
    }},
    ...
  ]
}}
"""
    
    user_prompt = f"""
Topic: {prompt}
Number of slides: {slide_count}
Language: {language}

Generate the JSON outline now.
"""

    try:
        logger.info(f"Generating outline for topic: {prompt}")
        response = client.chat.completions.create(
            model="gpt-4o-mini", 
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            response_format={"type": "json_object"}
        )
        
        response_content = response.choices[0].message.content
        logger.info(f"Raw AI response for outline: {response_content}")
        
        json_data = json.loads(response_content)
        
        # The model should return {"slides": [...]}
        if isinstance(json_data, dict) and "slides" in json_data and isinstance(json_data["slides"], list):
            return json_data["slides"]
        
        raise ValueError("The AI response did not contain a valid 'slides' array.")

    except json.JSONDecodeError as e:
        logger.error(f"Error decoding JSON from AI response: {e}")
        logger.error(f"Problematic response content: {response_content}")
        return None
    except Exception as e:
        logger.error(f"Error generating structured outline: {str(e)}")
        return None


def generate_presentation_config(content, template_structure):
    """Generates a presentation configuration using OpenAI GPT based on content and template structure."""
    prompt = f"""
    You are an assistant that creates PowerPoint presentations based on a template and content.

    Here is the structure of the PowerPoint template:
    {json.dumps(template_structure, indent=2)}

    The template has multiple slide masters, each with several layouts. Each layout contains placeholders of different types, such as Title, Body, Picture, etc., identified by idx, id, name, and type.

    You would be provided with the presentation content's outline.

    Create a JSON configuration for a PowerPoint presentation that uses the layouts from the template and fills the placeholders with appropriate content. The JSON should have the following format:

    {{
      "slides": [
        {{
          "master_idx": <master index>,
          "layout_idx": <layout index within the master>,
          "placeholders": {{
            "<placeholder_idx>": {{
              "type": "<placeholder_type>",
              "content": "<text or search query>"
            }},
            ...
          }}
        }},
        ...
      ]
    }}

    For placeholders of type Title, Body, Center Title, Subtitle, Date, Footer, Header, Slide Number, provide the text to insert in 'content'.

    For placeholders of type Picture, provide a search query for Pexels in 'content', which will be used to fetch an image from Pexels and insert it into the placeholder.

    Ensure that:
    - The presentation uses a variety of layouts appropriately (e.g., title slide, content slides with text and images).
    - The content is split across multiple slides if necessary (aim for an optimume numder of slides using different layouts based on the content outline).
    - Only text and picture placeholders are filled; other placeholders are ignored.
    - The placeholder_idx matches the idx in the template structure.
    - For picture placeholders, provide meaningful search queries based on the content.
    """

    try:
        response = client.chat.completions.create(
            model="gpt-4",
            messages=[
                {"role": "system", "content": prompt},
                {"role": "user", "content": content}
            ],
            timeout=60  # 60 second timeout for OpenAI API
        )
        return response.choices[0].message.content
    except Exception as e:
        logger.error(f"Error generating configuration: {str(e)}")
        return None

def generate_pptx_from_config(template_path, json_config):
    """Generates a PowerPoint file from the template and JSON configuration, including image insertion."""
    prs = Presentation(template_path)
    config = json.loads(json_config)

    for slide_data in config["slides"]:
        master_idx = slide_data["master_idx"]
        layout_idx = slide_data["layout_idx"]
        placeholders_data = slide_data["placeholders"]

        master = prs.slide_masters[master_idx]
        layout = master.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(layout)

        for ph_idx_str, ph_data in placeholders_data.items():
            ph_idx = int(ph_idx_str)
            placeholder = next((shape for shape in slide.placeholders if shape.placeholder_format.idx == ph_idx), None)
            if placeholder:
                placeholder_type = ph_data["type"]
                content = ph_data["content"]
                if placeholder_type in ["Title", "Body", "Center Title", "Subtitle", "Date", "Footer", "Header", "Slide Number"]:
                    placeholder.text = content
                elif placeholder_type == "Picture":
                    query = content
                    image_url = get_pexels_image_url(query, PEXELS_API_KEY)
                    if image_url:
                        try:
                            response = requests.get(image_url, stream=True, timeout=15)
                            if response.status_code == 200:
                                image_stream = io.BytesIO(response.content)
                                placeholder.insert_picture(image_stream)
                        except Exception as e:
                            logger.error(f"Error inserting image: {e}")

    pptx_bytes = io.BytesIO()
    # Delete the first slide (redundant slide that is not used in template)
    if len(prs.slides) > 0:
        slide_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(slide_id)
        del prs.slides._sldIdLst[0]

    
    prs.save(pptx_bytes)
    pptx_bytes.seek(0)
    return pptx_bytes

def get_libreoffice_path():
    """Determines the path to the LibreOffice executable."""
    if platform.system() == "Darwin":  # macOS
        return "/opt/homebrew/bin/libreoffice"
    elif platform.system() == "Linux":  # Ubuntu/Debian
        return "/usr/bin/libreoffice"
    else:  # Fallback for other systems
        return "libreoffice"

def check_libreoffice_installed():
    """Check if LibreOffice is installed and accessible."""
    libreoffice_path = get_libreoffice_path()
    try:
        if not os.path.exists(libreoffice_path) and libreoffice_path != "libreoffice":
             return False
        result = subprocess.run([libreoffice_path, "--version"], 
                              capture_output=True, text=True, timeout=10)
        return result.returncode == 0
    except (subprocess.TimeoutExpired, FileNotFoundError):
        return False

def convert_pptx_to_pdf(pptx_stream):
    """Converts a PowerPoint file to PDF using LibreOffice headless mode."""
    
    # Check if LibreOffice is installed
    if not check_libreoffice_installed():
        logger.error("LibreOffice is not installed or not accessible")
        raise Exception(
            "LibreOffice is required for PDF conversion. "
            "Please install LibreOffice:\n"
            "- macOS: brew install libreoffice\n"
            "- Ubuntu: sudo apt-get install libreoffice\n"
            "- Windows: Download from https://www.libreoffice.org/"
        )
    
    try:
        # Create temporary directory
        with tempfile.TemporaryDirectory() as temp_dir:
            # Save PPTX to temporary file
            pptx_path = os.path.join(temp_dir, "presentation.pptx")
            with open(pptx_path, "wb") as f:
                f.write(pptx_stream.getvalue())
            
            # Convert to PDF using LibreOffice
            logger.info("Converting PPTX to PDF using LibreOffice...")
            libreoffice_path = get_libreoffice_path()

            result = subprocess.run([
                libreoffice_path, "--headless", "--convert-to", "pdf", 
                "--outdir", temp_dir, pptx_path
            ], capture_output=True, text=True, timeout=60)
            
            if result.returncode != 0:
                logger.error(f"LibreOffice conversion failed: {result.stderr}")
                raise Exception(f"PDF conversion failed: {result.stderr}")
            
            # Read the generated PDF
            pdf_path = os.path.join(temp_dir, "presentation.pdf")
            if not os.path.exists(pdf_path):
                raise Exception("PDF file was not generated")
            
            with open(pdf_path, "rb") as pdf_file:
                pdf_content = pdf_file.read()
            
            logger.info("Successfully converted PPTX to PDF")
            return io.BytesIO(pdf_content)
            
    except subprocess.TimeoutExpired:
        logger.error("LibreOffice conversion timed out")
        raise Exception("PDF conversion timed out")
    except FileNotFoundError:
        logger.error("LibreOffice executable not found")
        raise Exception(
            "LibreOffice is not installed or not in PATH. "
            "Please install LibreOffice to enable PDF conversion."
        )
    except Exception as e:
        logger.error(f"Error converting PPTX to PDF: {e}")
        raise Exception(f"PDF conversion failed: {str(e)}")

class SlideRequest(BaseModel):
    content: str

class SlideRequestWithTemplate(BaseModel):
    title: str
    outlines: list
    templateId: str
    language: str = "english"

class ConversionRequest(BaseModel):
    file_url: str


class OutlineRequest(BaseModel):
    prompt: str
    slideCount: int
    language: str



app = FastAPI(title="Slide Generator API")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"], 
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)
@app.get("/template_structure")
def api_template_structure():
    try:
        struct = get_template_structure(DEFAULT_TEMPLATE_PATH)
        return struct
    except Exception as e:
        logger.error(f"Error in template_structure API: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/generate_outline")
def api_generate_outline(request: OutlineRequest):
    try:
        logger.info(f"Received outline generation request: {request.prompt}")
        
        outlines = generate_structured_outline(
            prompt=request.prompt,
            slide_count=request.slideCount,
            language=request.language
        )
        
        if not outlines:
            raise HTTPException(
                status_code=500, 
                detail="Failed to generate presentation outline from AI model."
            )
            
        logger.info(f"Successfully generated {len(outlines)} outlines.")
        
        return {"outlines": outlines}

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error in generate_outline API: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/generate_slide")
def api_generate_slide(request: SlideRequest):
    try:
        struct = get_template_structure(DEFAULT_TEMPLATE_PATH)
        outline = generate_content_outline(request.content)
        if not outline:
            raise HTTPException(status_code=400, detail="Failed to generate content outline")
        config = generate_presentation_config(outline, struct)
        if not config:
            raise HTTPException(status_code=400, detail="Failed to generate presentation config")
        pptx_io = generate_pptx_from_config(DEFAULT_TEMPLATE_PATH, config)
        return StreamingResponse(
            pptx_io,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": "attachment; filename=slidex_presentation.pptx"}
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error in generate_slide API: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/generate_slide_with_template")
def api_generate_slide_with_template(request: SlideRequestWithTemplate):
    try:
        # Debug: Log the received outlines
        logger.info(f"🔍 Received request for template: {request.templateId}")
        logger.info(f"📝 Total slides received: {len(request.outlines)}")
        for i, outline in enumerate(request.outlines):
            logger.info(f"📄 Slide {i + 1}: '{outline.get('title', 'Untitled')}' - {len(outline.get('content', []))} points")
            logger.info(f"   Content: {outline.get('content', [])}")
        
        # Get template path
        template_path = TEMPLATE_PATHS.get(request.templateId)
        if not template_path:
            raise HTTPException(status_code=400, detail=f"Template '{request.templateId}' not found")
        
        # Check if template file exists
        if not os.path.exists(template_path):
            raise HTTPException(status_code=400, detail=f"Template file '{template_path}' not found")
        
        # Get template structure
        struct = get_template_structure(template_path)
        
        # Convert outlines to content string
        content_parts = []
        for outline in request.outlines:
            content_parts.append(f"Slide: {outline.get('title', 'Untitled')}")
            for content_item in outline.get('content', []):
                content_parts.append(f"- {content_item}")
            content_parts.append("")  # Empty line between slides
        
        content_string = "\n".join(content_parts)
        logger.info(f"📋 Generated content string:\n{content_string}")
        
        # Generate presentation config using AI
        config = generate_presentation_config(content_string, struct)
        if not config:
            raise HTTPException(status_code=400, detail="Failed to generate presentation config")
        
        # Generate PPTX file
        pptx_io = generate_pptx_from_config(template_path, config)
        
        # Return as streaming response
        return StreamingResponse(
            pptx_io,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="{request.title}.pptx"'}
        )
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error in generate_slide_with_template API: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/convert_to_pdf")
def api_convert_to_pdf(request: ConversionRequest):
    try:
        response = requests.get(request.file_url)
        if response.status_code != 200:
            raise HTTPException(status_code=400, detail="Failed to fetch PPTX file")

        pptx_stream = io.BytesIO(response.content)
        pdf_io = convert_pptx_to_pdf(pptx_stream)
        
        return StreamingResponse(
            pdf_io,
            media_type="application/pdf",
            headers={"Content-Disposition": 'attachment; filename="presentation.pdf"'}
        )
            
    except Exception as e:
        logger.error(f"Error converting to PDF: {e}")
        raise HTTPException(status_code=500, detail=str(e))
